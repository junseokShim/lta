"""
문서 처리 도구
여러 형식의 문서를 읽고, 요약용 구조로 파싱하고, 프레젠테이션도 생성한다.
"""

import csv
import io
import json
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Optional

from ..logging_utils import get_logger

logger = get_logger("tools.document")


@dataclass
class ParsedDocument:
    """파싱된 문서 정보."""

    file_path: str
    file_type: str
    title: str
    content: str
    metadata: dict
    sections: list[dict]  # [{"title": "...", "content": "..."}]
    word_count: int
    char_count: int


class DocumentTool:
    """
    문서 읽기/쓰기/변환 도구.
    텍스트 기반 문서와 일부 바이너리 문서(PDF, DOCX)를 지원한다.
    """

    PARSERS = {
        ".txt": "_parse_text",
        ".md": "_parse_markdown",
        ".py": "_parse_code",
        ".js": "_parse_code",
        ".ts": "_parse_code",
        ".json": "_parse_json",
        ".yaml": "_parse_yaml",
        ".yml": "_parse_yaml",
        ".csv": "_parse_csv",
        ".html": "_parse_text",
        ".css": "_parse_code",
        ".sh": "_parse_code",
        ".toml": "_parse_text",
        ".pdf": "_parse_pdf",
        ".docx": "_parse_docx",
    }

    BINARY_PARSERS = {".pdf", ".docx"}

    def __init__(self, workspace_root: str):
        self.workspace_root = Path(workspace_root).resolve()

    def set_workspace_root(self, workspace_root: str) -> None:
        """작업 대상 루트를 동적으로 바꾼다."""
        self.workspace_root = Path(workspace_root).resolve()

    def read_document(self, file_path: str) -> ParsedDocument:
        """
        파일을 읽어 ParsedDocument 로 반환한다.
        """
        path = self._resolve_path(file_path)
        ext = path.suffix.lower()
        parser_name = self.PARSERS.get(ext, "_parse_text")
        parser = getattr(self, parser_name, self._parse_text)

        if ext in self.BINARY_PARSERS:
            doc = parser(str(path))
        else:
            try:
                raw_content = path.read_text(encoding="utf-8", errors="replace")
            except Exception as exc:
                raise IOError(f"파일 읽기 실패: {file_path} - {exc}") from exc
            doc = parser(raw_content, str(path))

        doc.file_path = str(path)
        doc.file_type = ext.lstrip(".")
        if not doc.char_count:
            doc.char_count = len(doc.content)
        if not doc.word_count:
            doc.word_count = len(doc.content.split())

        logger.debug("문서 파싱 완료: %s (%s words)", file_path, doc.word_count)
        return doc

    def write_document(
        self,
        file_path: str,
        content: str,
        doc_type: Optional[str] = None,
    ) -> str:
        """
        텍스트 문서를 저장한다.
        """
        _ = doc_type  # 현재는 확장자 기반 저장만 사용한다.
        path = self._resolve_path(file_path)
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(content, encoding="utf-8")
        logger.info("문서 저장: %s", file_path)
        return str(path)

    def create_presentation(
        self,
        title: str,
        sections: list[dict],
        output_path: str,
        subtitle: Optional[str] = None,
    ) -> str:
        """
        섹션 목록을 PPTX 파일로 만든다.
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError as exc:
            raise ImportError(
                "PPTX 생성에는 'python-pptx' 패키지가 필요합니다. "
                "requirements.txt 를 다시 설치해 주세요."
            ) from exc

        normalized_sections = self._normalize_sections_for_slides(sections)
        path = self._resolve_path(output_path)
        path.parent.mkdir(parents=True, exist_ok=True)

        presentation = Presentation()

        title_layout = presentation.slide_layouts[0]
        title_slide = presentation.slides.add_slide(title_layout)
        if title_slide.shapes.title:
            title_slide.shapes.title.text = title
        if len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = subtitle or ""

        content_layout = (
            presentation.slide_layouts[1]
            if len(presentation.slide_layouts) > 1
            else presentation.slide_layouts[0]
        )

        for section in normalized_sections:
            slide = presentation.slides.add_slide(content_layout)
            slide_title = section.get("title") or "Slide"
            if slide.shapes.title:
                slide.shapes.title.text = slide_title

            bullets = section.get("bullets") or self._content_to_bullets(section.get("content", ""))
            text_frame = None

            if len(slide.placeholders) > 1 and hasattr(slide.placeholders[1], "text_frame"):
                text_frame = slide.placeholders[1].text_frame
            else:
                textbox = slide.shapes.add_textbox(
                    Inches(1.0),
                    Inches(1.8),
                    Inches(8.0),
                    Inches(4.5),
                )
                text_frame = textbox.text_frame

            text_frame.clear()
            if bullets:
                for index, bullet in enumerate(bullets):
                    paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
                    paragraph.text = bullet
                    paragraph.level = 0
            else:
                paragraph = text_frame.paragraphs[0]
                paragraph.text = (section.get("content", "") or "").strip() or "내용 없음"

        presentation.save(str(path))
        logger.info("프레젠테이션 저장: %s", path)
        return str(path)

    def create_presentation_from_file(
        self,
        source_path: str,
        output_path: str,
        title: Optional[str] = None,
        subtitle: Optional[str] = None,
    ) -> str:
        """
        문서를 읽어 PPTX 로 변환한다.
        """
        doc = self.read_document(source_path)
        deck_title = title or doc.title or Path(source_path).stem
        sections = doc.sections or [{"title": deck_title, "content": doc.content}]
        return self.create_presentation(deck_title, sections, output_path, subtitle=subtitle)

    def generate_report(
        self,
        title: str,
        sections: list[dict],
        output_path: str,
        format: str = "markdown",
    ) -> str:
        """
        보고서를 생성한다.
        """
        if format == "markdown":
            content = self._build_markdown_report(title, sections)
            ext = ".md"
        elif format == "html":
            content = self._build_html_report(title, sections)
            ext = ".html"
        else:
            content = self._build_text_report(title, sections)
            ext = ".txt"

        if not output_path.endswith(ext):
            output_path = output_path + ext if "." not in Path(output_path).name else output_path

        return self.write_document(output_path, content)

    def summarize_document(self, doc: ParsedDocument, max_chars: int = 2000) -> str:
        """
        문서 요약. LLM 없이 핵심 섹션 위주로 압축한다.
        """
        if len(doc.content) <= max_chars:
            return doc.content

        if doc.sections:
            summary_parts = [f"**{doc.title}**\n"]
            remaining = max_chars - len(doc.title)
            per_section = max(remaining // max(len(doc.sections), 1), 200)

            for section in doc.sections[:5]:
                title = section.get("title", "")
                content = section.get("content", "")[:per_section]
                if title:
                    summary_parts.append(f"\n## {title}\n{content}...")
                else:
                    summary_parts.append(f"\n{content}...")

            return "\n".join(summary_parts)

        head = doc.content[: max_chars // 2]
        tail = doc.content[-max_chars // 4 :]
        return f"{head}\n\n[...중략...]\n\n{tail}"

    def extract_code_blocks(self, content: str) -> list[dict]:
        """
        마크다운에서 코드 블록을 추출한다.
        """
        blocks = []
        pattern = r"```(\w*)\n(.*?)```"
        for match in re.finditer(pattern, content, re.DOTALL):
            blocks.append(
                {
                    "language": match.group(1) or "text",
                    "code": match.group(2).strip(),
                }
            )
        return blocks

    def merge_documents(self, docs: list[ParsedDocument], output_path: str) -> str:
        """여러 문서를 하나의 리포트로 합친다."""
        sections = []
        for doc in docs:
            sections.append(
                {
                    "title": doc.title or Path(doc.file_path).name,
                    "content": doc.content,
                }
            )
        return self.generate_report("통합 문서", sections, output_path)

    def _parse_text(self, content: str, file_path: str) -> ParsedDocument:
        """일반 텍스트 파싱."""
        lines = content.splitlines()
        title = lines[0].strip() if lines else Path(file_path).stem
        return ParsedDocument(
            file_path=file_path,
            file_type="txt",
            title=title,
            content=content,
            metadata={},
            sections=[{"title": "", "content": content}],
            word_count=0,
            char_count=0,
        )

    def _parse_markdown(self, content: str, file_path: str) -> ParsedDocument:
        """마크다운을 H2 기준으로 나눈다."""
        lines = content.splitlines()
        title = Path(file_path).stem

        for line in lines:
            if line.startswith("# "):
                title = line[2:].strip()
                break

        sections = []
        current_section = {"title": "", "content": []}

        for line in lines:
            if line.startswith("## "):
                if current_section["content"]:
                    sections.append(
                        {
                            "title": current_section["title"],
                            "content": "\n".join(current_section["content"]).strip(),
                        }
                    )
                current_section = {"title": line[3:].strip(), "content": []}
            else:
                current_section["content"].append(line)

        if current_section["content"]:
            sections.append(
                {
                    "title": current_section["title"],
                    "content": "\n".join(current_section["content"]).strip(),
                }
            )

        return ParsedDocument(
            file_path=file_path,
            file_type="md",
            title=title,
            content=content,
            metadata={},
            sections=sections if sections else [{"title": "", "content": content}],
            word_count=0,
            char_count=0,
        )

    def _parse_code(self, content: str, file_path: str) -> ParsedDocument:
        """코드 파일을 간단히 구조화한다."""
        path = Path(file_path)
        ext = path.suffix.lower()
        sections = []

        if ext == ".py":
            for match in re.finditer(r"^(class|def)\s+(\w+)", content, re.MULTILINE):
                sections.append({"title": f"{match.group(1)} {match.group(2)}", "content": ""})

        return ParsedDocument(
            file_path=file_path,
            file_type=ext.lstrip("."),
            title=path.stem,
            content=content,
            metadata={"language": ext.lstrip(".")},
            sections=sections if sections else [{"title": "", "content": content}],
            word_count=0,
            char_count=0,
        )

    def _parse_json(self, content: str, file_path: str) -> ParsedDocument:
        """JSON 파일 파싱."""
        try:
            data = json.loads(content)
            pretty = json.dumps(data, ensure_ascii=False, indent=2)
            metadata = {
                "keys": list(data.keys()) if isinstance(data, dict) else [],
                "type": type(data).__name__,
            }
        except json.JSONDecodeError as exc:
            pretty = content
            metadata = {"error": str(exc)}

        return ParsedDocument(
            file_path=file_path,
            file_type="json",
            title=Path(file_path).stem,
            content=pretty,
            metadata=metadata,
            sections=[{"title": "", "content": pretty}],
            word_count=0,
            char_count=0,
        )

    def _parse_yaml(self, content: str, file_path: str) -> ParsedDocument:
        """YAML 파일 파싱."""
        try:
            import yaml

            data = yaml.safe_load(content)
            metadata = {"keys": list(data.keys()) if isinstance(data, dict) else []}
        except Exception as exc:
            metadata = {"error": str(exc)}

        return ParsedDocument(
            file_path=file_path,
            file_type="yaml",
            title=Path(file_path).stem,
            content=content,
            metadata=metadata,
            sections=[{"title": "", "content": content}],
            word_count=0,
            char_count=0,
        )

    def _parse_csv(self, content: str, file_path: str) -> ParsedDocument:
        """CSV 파일 파싱."""
        try:
            reader = csv.DictReader(io.StringIO(content))
            rows = list(reader)
            headers = reader.fieldnames or []
            preview = f"헤더: {', '.join(headers)}\n행 수: {len(rows)}\n\n"
            if rows:
                preview += "첫 5행\n"
                for row in rows[:5]:
                    preview += str(dict(row)) + "\n"
            metadata = {"headers": list(headers), "row_count": len(rows)}
        except Exception as exc:
            preview = content
            metadata = {"error": str(exc)}

        return ParsedDocument(
            file_path=file_path,
            file_type="csv",
            title=Path(file_path).stem,
            content=content,
            metadata=metadata,
            sections=[{"title": "데이터 미리보기", "content": preview}],
            word_count=0,
            char_count=0,
        )

    def _parse_pdf(self, file_path: str) -> ParsedDocument:
        """PDF 텍스트 추출."""
        try:
            from PyPDF2 import PdfReader
        except ImportError as exc:
            raise ImportError(
                "PDF 읽기에는 'PyPDF2' 패키지가 필요합니다. requirements.txt 를 다시 설치해 주세요."
            ) from exc

        with open(file_path, "rb") as handle:
            reader = PdfReader(handle)
            pages = []
            sections = []

            for page_index, page in enumerate(reader.pages, start=1):
                text = (page.extract_text() or "").strip()
                if not text:
                    continue
                pages.append(text)
                sections.append({"title": f"Page {page_index}", "content": text})

            metadata = {"page_count": len(reader.pages)}
            if reader.metadata:
                metadata["pdf_metadata"] = {
                    str(key).lstrip("/"): str(value)
                    for key, value in dict(reader.metadata).items()
                    if value is not None
                }

        content = "\n\n".join(pages)
        title = (
            metadata.get("pdf_metadata", {}).get("Title")
            or metadata.get("pdf_metadata", {}).get("title")
            or Path(file_path).stem
        )

        return ParsedDocument(
            file_path=file_path,
            file_type="pdf",
            title=title,
            content=content,
            metadata=metadata,
            sections=sections if sections else [{"title": "", "content": content}],
            word_count=len(content.split()),
            char_count=len(content),
        )

    def _parse_docx(self, file_path: str) -> ParsedDocument:
        """DOCX 텍스트 추출."""
        try:
            from docx import Document as DocxDocument
        except ImportError as exc:
            raise ImportError(
                "DOCX 읽기에는 'python-docx' 패키지가 필요합니다. requirements.txt 를 다시 설치해 주세요."
            ) from exc

        doc = DocxDocument(file_path)
        text_parts = []
        sections = []

        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            text_parts.append(text)
            if getattr(paragraph.style, "name", "").startswith("Heading"):
                sections.append({"title": text, "content": ""})
            elif sections:
                current = sections[-1]
                current["content"] = (current["content"] + "\n" + text).strip()

        text = "\n".join(text_parts)
        if not sections:
            sections = [{"title": "", "content": text}]

        title = text_parts[0] if text_parts else Path(file_path).stem
        return ParsedDocument(
            file_path=file_path,
            file_type="docx",
            title=title,
            content=text,
            metadata={},
            sections=sections,
            word_count=len(text.split()),
            char_count=len(text),
        )

    def _resolve_path(self, file_path: str) -> Path:
        """절대/상대 경로를 현재 루트 기준으로 해석한다."""
        path = Path(file_path)
        if path.is_absolute():
            return path
        return (self.workspace_root / file_path).resolve()

    def _content_to_bullets(self, content: str, max_bullets: int = 6) -> list[str]:
        """문단/목록 텍스트를 슬라이드용 bullet 리스트로 바꾼다."""
        bullets = []
        for raw_line in content.splitlines():
            line = raw_line.strip()
            if not line:
                continue
            line = re.sub(r"^[-*•]\s+", "", line)
            line = re.sub(r"^\d+[.)]\s+", "", line)
            if len(line) > 180:
                line = line[:177] + "..."
            bullets.append(line)

        if not bullets and content.strip():
            normalized = re.sub(r"\s+", " ", content).strip()
            chunks = re.split(r"(?<=[.!?])\s+", normalized)
            bullets = [chunk.strip() for chunk in chunks if chunk.strip()]

        return bullets[:max_bullets]

    def _normalize_sections_for_slides(self, sections: list[dict], max_bullets: int = 6) -> list[dict]:
        """슬라이드 생성을 위해 섹션 목록을 정규화한다."""
        normalized = []

        for index, section in enumerate(sections, start=1):
            title = (section.get("title") or "").strip() or f"Slide {index}"
            content = (section.get("content") or "").strip()
            bullets = section.get("bullets") or self._content_to_bullets(content, max_bullets=max_bullets * 2)

            if bullets:
                for chunk_index in range(0, len(bullets), max_bullets):
                    chunk = bullets[chunk_index : chunk_index + max_bullets]
                    suffix = "" if chunk_index == 0 else f" ({chunk_index // max_bullets + 1})"
                    normalized.append({"title": f"{title}{suffix}", "bullets": chunk, "content": content})
            else:
                normalized.append({"title": title, "bullets": [], "content": content})

        return normalized or [{"title": "Slide 1", "bullets": [], "content": ""}]

    def _build_markdown_report(self, title: str, sections: list[dict]) -> str:
        """마크다운 리포트를 만든다."""
        from datetime import datetime

        lines = [
            f"# {title}",
            f"\n> 생성일: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n",
        ]
        for section in sections:
            section_title = section.get("title", "")
            content = section.get("content", "")
            if section_title:
                lines.append(f"\n## {section_title}\n")
            lines.append(content)
        return "\n".join(lines)

    def _build_html_report(self, title: str, sections: list[dict]) -> str:
        """HTML 리포트를 만든다."""
        from datetime import datetime
        import html as html_module

        body_parts = []
        for section in sections:
            section_title = section.get("title", "")
            content = html_module.escape(section.get("content", ""))
            if section_title:
                body_parts.append(f"<h2>{html_module.escape(section_title)}</h2>")
            body_parts.append(f"<pre>{content}</pre>")

        return f"""<!DOCTYPE html>
<html lang="ko">
<head><meta charset="UTF-8"><title>{html_module.escape(title)}</title>
<style>body{{font-family:sans-serif;max-width:900px;margin:0 auto;padding:20px}}
pre{{background:#f4f4f4;padding:15px;border-radius:5px;overflow:auto}}</style>
</head>
<body>
<h1>{html_module.escape(title)}</h1>
<p><em>생성일: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}</em></p>
{"".join(body_parts)}
</body></html>"""

    def _build_text_report(self, title: str, sections: list[dict]) -> str:
        """텍스트 리포트를 만든다."""
        from datetime import datetime

        lines = [
            "=" * 60,
            title.center(60),
            f"생성일: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}".center(60),
            "=" * 60,
        ]
        for section in sections:
            section_title = section.get("title", "")
            content = section.get("content", "")
            if section_title:
                lines.extend(["", "-" * 40, section_title, "-" * 40])
            lines.append(content)
        return "\n".join(lines)
